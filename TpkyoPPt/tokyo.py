from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Regiunea Industrială Tokyo"
subtitle.text = "Prezentare pentru titularizare\n[Mitroi Iuliana]\n[31.05.2024]"

# Slide 1: Localizare Geografică
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Localizare Geografică"
content = slide.placeholders[1]
content.text = "• Situată în estul Japoniei\n• Include orașele Tokyo, Yokohama, Kawasaki\n• Centrul economic al Japoniei"

# Slide 2: Istorie
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Istorie"
content = slide.placeholders[1]
content.text = "• Dezvoltare rapidă în perioada Meiji (1868-1912)\n• Reconstrucția post-Al Doilea Război Mondial\n• Creștere economică în perioada de după 1950"

# Slide 3: Industrii Principale
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Industrii Principale"
content = slide.placeholders[1]
content.text = "• Tehnologie și electronică\n• Industria auto\n• Finanțe și servicii bancare\n• Industria chimică\n• Producție și robotică"

# Slide 4: Infrastructură
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Infrastructură"
content = slide.placeholders[1]
content.text = "• Rețea extinsă de transport public: trenuri, metrou, autobuze\n• Aeroporturi internaționale: Narita și Haneda\n• Porturi maritime majore\n• Centre de expoziții și convenții"

# Slide 5: Importanța Economică
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Importanța Economică"
content = slide.placeholders[1]
content.text = "• Contribuție majoră la PIB-ul Japoniei\n• Hub pentru investiții internaționale\n• Inovație și cercetare\n• Centrul financiar al Asiei"

# Slide 6: Provocări și Perspective
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Provocări și Perspective"
content = slide.placeholders[1]
content.text = "• Probleme demografice: îmbătrânirea populației\n• Dezastre naturale: risc seismic\n• Tranziția către o economie verde\n• Adoptarea noilor tehnologii și digitalizarea"

# Slide 7: Concluzii
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Concluzii"
content = slide.placeholders[1]
content.text = "• Tokyo este un centru industrial și economic global\n• Importanța sa istorică și contemporană în Japonia\n• Provocările actuale și oportunitățile de viitor"

# Slide 8: Întrebări și Răspunsuri
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Întrebări și Răspunsuri"
content = slide.placeholders[1]
content.text = "• Spațiu pentru discuții și întrebări din partea audienței"

# Save the presentation
file_path = "Tokyo_Industrial_Region_Presentation.pptx"
prs.save(file_path)

print(f"Prezentarea a fost generată și salvată la: {file_path}")
